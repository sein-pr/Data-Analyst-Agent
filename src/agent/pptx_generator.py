from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
import io
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .analysis_engine import AnalysisResult
from .data_healer import MappingResult
from .brand import BrandGuidelines
from .logger import get_logger
from .quickchart_client import QuickChartClient

logger = get_logger(__name__)


@dataclass
class SlideContent:
    title: str
    bullets: List[str]


class PPTXGenerator:
    def __init__(
        self,
        brand: BrandGuidelines,
        logo_full_path: Path | None = None,
        logo_symbol_path: Path | None = None,
    ) -> None:
        self.brand = brand
        self.logo_full_path = logo_full_path
        self.logo_symbol_path = logo_symbol_path
        self.report_source = ""
        self.report_date = ""
        self.primary_font = "Calibri"
        self.quickchart = QuickChartClient()
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        self.margin = Inches(0.5)

    def build(
        self,
        analysis: AnalysisResult,
        output_path: Path,
        bullets: List[str],
        mapping: MappingResult,
        report_source: str,
        primary_font: str | None = None,
        department_label: str | None = None,
        previous_analysis: dict | None = None,
    ) -> Path:
        if primary_font:
            self.primary_font = primary_font
        self.report_source = report_source
        self.report_date = datetime.utcnow().strftime("%Y-%m-%d")
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        self._add_title_slide(prs, analysis, department_label)
        self._add_kpi_summary_slide(prs, bullets)
        self._add_kpi_slide(prs, analysis)
        if previous_analysis:
            self._add_comparison_slide(prs, analysis, previous_analysis)
        self._add_self_healing_slide(prs, analysis, mapping)
        if self._should_add_mapping_slide(mapping):
            self._add_mapping_detail_slide(prs, mapping)
        if analysis.monthly_revenue:
            self._add_mom_trend_slide(prs, analysis)
        self._add_data_quality_slide(prs, analysis)
        self._add_recommendations_slide(prs, analysis, bullets)
        self._apply_theme(prs)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        logger.info("Saved presentation to %s", output_path)
        return output_path

    def _apply_theme(self, prs: Presentation) -> None:
        for slide in prs.slides:
            self._set_slide_background(slide)
            self._add_footer_logo(slide)
            self._add_footer_meta(slide)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.primary_font
                        try:
                            if run.font.color.rgb is None:
                                run.font.color.rgb = RGBColor.from_string(
                                    self.brand.palette.primary[1:]
                                )
                        except Exception:  # noqa: BLE001
                            run.font.color.rgb = RGBColor.from_string(
                                self.brand.palette.primary[1:]
                            )

    def _set_slide_background(self, slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(self.brand.palette.neutral[1:])

    def _add_footer_logo(self, slide) -> None:
        if not self.logo_symbol_path or not self.logo_symbol_path.exists():
            return
        slide.shapes.add_picture(
            str(self.logo_symbol_path),
            Inches(8.6),
            Inches(6.95),
            width=Inches(0.55),
        )

    def _add_footer_meta(self, slide) -> None:
        if not self.report_date:
            return
        meta_box = slide.shapes.add_textbox(self.margin, Inches(7.05), Inches(7.5), Inches(0.35))
        tf = meta_box.text_frame
        tf.text = f"{self.report_date} | Source: {self.report_source}"
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor.from_string(self.brand.palette.primary[1:])

    def _add_header_bar(self, slide) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            self.slide_width,
            Inches(0.35),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(self.brand.palette.primary[1:])
        shape.line.fill.background()

    def _apply_font(self, text_frame, size: int | None = None, bold: bool | None = None) -> None:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.primary_font
                if size:
                    run.font.size = Pt(size)
                if bold is not None:
                    run.font.bold = bold

    def _shrink_text_to_fit(self, text_frame, max_chars: int, base_size: int, min_size: int) -> None:
        text = "".join(p.text for p in text_frame.paragraphs)
        if not text:
            return
        size = base_size
        if len(text) > max_chars:
            size = max(min_size, int(base_size - (len(text) - max_chars) / 40))
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)

    def _add_title_slide(self, prs: Presentation, analysis: AnalysisResult, department_label: str | None) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)

        if self.logo_full_path and self.logo_full_path.exists():
            slide.shapes.add_picture(
                str(self.logo_full_path),
                self.margin,
                Inches(0.45),
                width=Inches(2.8),
            )

        title_box = slide.shapes.add_textbox(self.margin, Inches(2.0), Inches(9.0), Inches(1.0))
        title_tf = title_box.text_frame
        dept = f"{department_label.title()} " if department_label else ""
        title_tf.text = f"{self.brand.name} {dept}Executive Summary"
        title_tf.paragraphs[0].font.size = Pt(30)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=30, bold=True)

        subtitle_box = slide.shapes.add_textbox(self.margin, Inches(2.9), Inches(8.5), Inches(0.7))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = self.brand.tagline or "Automated Insights"
        subtitle_tf.paragraphs[0].font.size = Pt(16)
        subtitle_tf.paragraphs[0].font.color.rgb = RGBColor.from_string(self.brand.palette.secondary[1:])
        self._apply_font(subtitle_tf, size=16)

    def _add_kpi_summary_slide(self, prs: Presentation, bullets: List[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)

        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(8.8), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "Executive Highlights"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        left = self.margin
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(5.2)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        bullets = [b for b in bullets if not b.lower().startswith("rec:")]
        if not bullets:
            bullets = ["Key performance is stable with no major anomalies detected."]
        tf.text = f"- {self._wrap_text(bullets[0], 70)}"
        tf.paragraphs[0].font.size = Pt(18)
        self._apply_font(tf, size=18)
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = f"- {self._wrap_text(bullet, 70)}"
            p.level = 0
            p.font.size = Pt(18)
        self._shrink_text_to_fit(tf, max_chars=700, base_size=18, min_size=14)

    def _add_kpi_slide(self, prs: Presentation, analysis: AnalysisResult) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)
        self._add_data_health_badge(slide, analysis)

        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(6.0), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "KPI Dashboard"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        kpi_table = slide.shapes.add_table(
            rows=max(2, len(analysis.kpis) + 1),
            cols=2,
            left=self.margin,
            top=Inches(1.5),
            width=Inches(4.6),
            height=Inches(3.0),
        ).table
        kpi_table.cell(0, 0).text = "KPI"
        kpi_table.cell(0, 1).text = "Value"
        self._style_table_header(kpi_table)
        row = 1
        kpi_font_size = 18 if len(analysis.kpis) > 6 else 20
        for key, value in analysis.kpis.items():
            kpi_table.cell(row, 0).text = key
            kpi_table.cell(row, 1).text = value
            self._apply_font(kpi_table.cell(row, 0).text_frame, size=kpi_font_size)
            self._apply_font(kpi_table.cell(row, 1).text_frame, size=kpi_font_size)
            row += 1

        if analysis.top_products:
            top_table = slide.shapes.add_table(
                rows=len(analysis.top_products) + 1,
                cols=2,
                left=Inches(5.3),
                top=Inches(1.5),
                width=Inches(4.2),
                height=Inches(2.4),
            ).table
            top_table.cell(0, 0).text = "Top Products"
            top_table.cell(0, 1).text = "Revenue"
            self._style_table_header(top_table)
            for idx, item in enumerate(analysis.top_products, start=1):
                category = item.get("Product Category", "")
                if len(category) > 28:
                    category = category[:25].rstrip() + "..."
                top_table.cell(idx, 0).text = category
                top_table.cell(idx, 1).text = item.get("Revenue", "")
                self._apply_font(top_table.cell(idx, 0).text_frame, size=16)
                self._apply_font(top_table.cell(idx, 1).text_frame, size=16)
        self._add_outliers_table(slide, analysis)
        self._add_top_products_chart(slide, analysis)

    def _add_comparison_slide(
        self, prs: Presentation, analysis: AnalysisResult, previous_analysis: dict
    ) -> None:
        prev_kpis = previous_analysis.get("kpis") or {}
        if not prev_kpis:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)

        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(8.8), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "Comparative Analysis"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        common_keys = [k for k in analysis.kpis.keys() if k in prev_kpis]
        if not common_keys:
            common_keys = list(analysis.kpis.keys())[:6]

        rows = min(len(common_keys) + 1, 7)
        table = slide.shapes.add_table(
            rows=rows,
            cols=4,
            left=self.margin,
            top=Inches(1.5),
            width=Inches(9.0),
            height=Inches(3.0),
        ).table
        table.cell(0, 0).text = "KPI"
        table.cell(0, 1).text = "Current"
        table.cell(0, 2).text = "Previous"
        table.cell(0, 3).text = "Delta"
        self._style_table_header(table)

        for idx, key in enumerate(common_keys[:6], start=1):
            current = analysis.kpis.get(key, "")
            previous = prev_kpis.get(key, "")
            delta = self._format_delta(current, previous)
            table.cell(idx, 0).text = key
            table.cell(idx, 1).text = str(current)
            table.cell(idx, 2).text = str(previous)
            table.cell(idx, 3).text = delta
            for col in range(4):
                self._apply_font(table.cell(idx, col).text_frame, size=16)

        note_box = slide.shapes.add_textbox(self.margin, Inches(4.9), Inches(9.0), Inches(1.5))
        tf = note_box.text_frame
        tf.text = "Variance is computed from the latest stored analysis for this schema."
        tf.paragraphs[0].font.size = Pt(12)
        self._apply_font(tf, size=12)

    def _add_self_healing_slide(
        self, prs: Presentation, analysis: AnalysisResult, mapping: MappingResult
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)

        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(8.8), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "Self-Healing Report"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        self._add_schema_summary_table(slide, mapping)

        body_box = slide.shapes.add_textbox(self.margin, Inches(1.5), Inches(6.2), Inches(4.8))
        body_tf = body_box.text_frame
        body_tf.text = "Schema alignment completed."
        body_tf.paragraphs[0].font.size = Pt(20)
        self._apply_font(body_tf, size=20)

        if mapping.mapping:
            mapping_title = body_tf.add_paragraph()
            mapping_title.text = "Mapped Columns:"
            mapping_title.font.size = Pt(18)
            mapping_title.font.bold = True
            for raw, target in mapping.mapping.items():
                p = body_tf.add_paragraph()
                p.text = f"{raw} -> {target}"
                p.level = 1
                p.font.size = Pt(18)

        if mapping.new_columns:
            new_title = body_tf.add_paragraph()
            new_title.text = "New Columns:"
            new_title.font.size = Pt(18)
            new_title.font.bold = True
            for col in mapping.new_columns:
                p = body_tf.add_paragraph()
                p.text = col
                p.level = 1
                p.font.size = Pt(18)

        if mapping.unmapped_required:
            missing_title = body_tf.add_paragraph()
            missing_title.text = "Missing Required Columns:"
            missing_title.font.size = Pt(18)
            missing_title.font.bold = True
            for col in mapping.unmapped_required:
                p = body_tf.add_paragraph()
                p.text = col
                p.level = 1
                p.font.size = Pt(18)
        self._shrink_text_to_fit(body_tf, max_chars=600, base_size=18, min_size=14)

    def _add_schema_summary_table(self, slide, mapping: MappingResult) -> None:
        table = slide.shapes.add_table(
            rows=4,
            cols=2,
            left=Inches(6.9),
            top=Inches(1.5),
            width=Inches(2.6),
            height=Inches(2.2),
        ).table
        table.cell(0, 0).text = "Schema Metric"
        table.cell(0, 1).text = "Count"
        self._style_table_header(table)

        metrics = [
            ("Mapped Columns", str(len(mapping.mapping))),
            ("New Columns", str(len(mapping.new_columns))),
            ("Missing Required", str(len(mapping.unmapped_required))),
        ]
        for idx, (label, value) in enumerate(metrics, start=1):
            table.cell(idx, 0).text = label
            table.cell(idx, 1).text = value

    def _add_recommendations_slide(
        self, prs: Presentation, analysis: AnalysisResult, bullets: List[str]
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)

        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(8.8), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "Recommendations"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        left = self.margin
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(5.0)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        recs = self._extract_recommendations(bullets)
        if not recs:
            recs = ["Review top-performing categories and investigate outliers."]
        tf.text = f"- {self._wrap_text(recs[0], 70)}"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        self._apply_font(tf, size=22)
        for bullet in recs[1:]:
            p = tf.add_paragraph()
            p.text = f"- {self._wrap_text(bullet, 70)}"
            p.level = 0
            p.font.size = Pt(22)

        self._add_slide_notes(slide, f"Recommendations: {len(recs)}")
        self._shrink_text_to_fit(tf, max_chars=700, base_size=22, min_size=16)

    def _add_data_quality_slide(self, prs: Presentation, analysis: AnalysisResult) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)

        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(8.8), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "Dataset Overview"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        table = slide.shapes.add_table(
            rows=6,
            cols=2,
            left=self.margin,
            top=Inches(1.5),
            width=Inches(4.6),
            height=Inches(3.0),
        ).table
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        self._style_table_header(table)

        metrics = [
            ("Rows", analysis.data_quality.get("rows", "0")),
            ("Columns", analysis.data_quality.get("columns", "0")),
            ("Missing Cells", analysis.data_quality.get("missing_cells", "0")),
            ("Missing %", analysis.data_quality.get("missing_pct", "0.00%")),
            ("Duplicate Rows", analysis.data_quality.get("duplicate_rows", "0")),
        ]
        for idx, (label, value) in enumerate(metrics, start=1):
            table.cell(idx, 0).text = label
            table.cell(idx, 1).text = value
            self._apply_font(table.cell(idx, 0).text_frame, size=18)
            self._apply_font(table.cell(idx, 1).text_frame, size=18)

        note_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(4.2), Inches(3.0))
        tf = note_box.text_frame
        tf.text = "Dataset summary and quality metrics."
        tf.paragraphs[0].font.size = Pt(18)
        self._apply_font(tf, size=18)

        schema_box = slide.shapes.add_textbox(self.margin, Inches(4.7), Inches(9.0), Inches(1.3))
        schema_tf = schema_box.text_frame
        schema_tf.text = self._wrap_text(
            f"Columns: {analysis.schema_overview.get('columns', '')}", 90
        )
        schema_tf.paragraphs[0].font.size = Pt(12)
        self._apply_font(schema_tf, size=12)

        self._add_slide_notes(slide, f"Missing cell %: {analysis.data_quality.get('missing_pct', '')}")

    def _style_table_header(self, table) -> None:
        for col_idx in range(len(table.columns)):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(self.brand.palette.primary[1:])
            if cell.text_frame.paragraphs:
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.name = self.primary_font

    def _add_top_products_chart(self, slide, analysis: AnalysisResult) -> None:
        if not analysis.top_products:
            return
        labels = []
        for item in analysis.top_products:
            label = item.get("Product Category", "")
            labels.append(label[:12] + "..." if len(label) > 12 else label)
        values = []
        for item in analysis.top_products:
            raw = item.get("Revenue", "0").replace(",", "")
            try:
                values.append(float(raw))
            except ValueError:
                values.append(0.0)

        chart_config = {
            "type": "bar",
            "data": {
                "labels": labels,
                "datasets": [
                    {
                        "label": "Revenue",
                        "data": values,
                        "backgroundColor": self.brand.palette.secondary,
                        "borderColor": self.brand.palette.primary,
                        "borderWidth": 1,
                    }
                ],
            },
            "options": {
                "plugins": {"legend": {"display": False}},
                "scales": {"x": {"ticks": {"autoSkip": True}}},
            },
        }

        image = self.quickchart.render_chart(chart_config, width=900, height=160, background=self.brand.palette.neutral)
        if image:
            slide.shapes.add_picture(
                io.BytesIO(image),
                self.margin,
                Inches(4.8),
                width=Inches(9.0),
                height=Inches(1.4),
            )
            self._add_slide_notes(slide, "Top products chart rendered (QuickChart).")
            return

        # Fallback to built-in chart if QuickChart fails
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
        except Exception:  # noqa: BLE001
            return

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("Revenue", values)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            self.margin,
            Inches(4.8),
            Inches(9.0),
            Inches(1.4),
            chart_data,
        ).chart
        chart.has_legend = False
        self._add_slide_notes(slide, "Top products chart rendered.")

    def _add_mom_trend_slide(self, prs: Presentation, analysis: AnalysisResult) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)
        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(8.8), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "Revenue Trend (Monthly)"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        labels = [item["Month"] for item in analysis.monthly_revenue]
        values = []
        for item in analysis.monthly_revenue:
            raw = item["Revenue"].replace(",", "")
            try:
                values.append(float(raw))
            except ValueError:
                values.append(0.0)

        chart_config = {
            "type": "line",
            "data": {
                "labels": labels,
                "datasets": [
                    {
                        "label": "Revenue",
                        "data": values,
                        "borderColor": self.brand.palette.primary,
                        "backgroundColor": self.brand.palette.secondary,
                        "fill": False,
                        "tension": 0.2,
                    }
                ],
            },
            "options": {"plugins": {"legend": {"display": False}}},
        }

        image = self.quickchart.render_chart(chart_config, width=900, height=420, background=self.brand.palette.neutral)
        if image:
            slide.shapes.add_picture(
                io.BytesIO(image),
                self.margin,
                Inches(1.5),
                width=Inches(9.0),
                height=Inches(4.6),
            )
            self._add_slide_notes(slide, f"Monthly revenue points: {len(analysis.monthly_revenue)} (QuickChart)")
            return

        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
        except Exception:  # noqa: BLE001
            return

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("Revenue", values)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            self.margin,
            Inches(1.5),
            Inches(9.0),
            Inches(4.6),
            chart_data,
        ).chart
        chart.has_legend = False
        self._add_slide_notes(slide, f"Monthly revenue points: {len(analysis.monthly_revenue)}")

    def _add_outliers_table(self, slide, analysis: AnalysisResult) -> None:
        if not analysis.outliers:
            return
        table = slide.shapes.add_table(
            rows=len(analysis.outliers) + 1,
            cols=2,
            left=Inches(5.3),
            top=Inches(4.0),
            width=Inches(4.2),
            height=Inches(1.0),
        ).table
        table.cell(0, 0).text = "Outlier"
        table.cell(0, 1).text = "Revenue"
        self._style_table_header(table)
        for idx, item in enumerate(analysis.outliers, start=1):
            label = item.get("Product Category", "") or "Unspecified"
            if len(label) > 22:
                label = label[:19].rstrip() + "..."
            table.cell(idx, 0).text = label
            table.cell(idx, 1).text = item.get("Revenue", "")
            self._apply_font(table.cell(idx, 0).text_frame, size=14)
            self._apply_font(table.cell(idx, 1).text_frame, size=14)

    def _add_data_health_badge(self, slide, analysis: AnalysisResult) -> None:
        status = "Green"
        color = "#4CAF50"
        if analysis.summary.get("outlier_count", "0") != "0":
            status = "Amber"
            color = "#F4A261"
        if analysis.summary.get("kpi_count", "0") == "0":
            status = "Red"
            color = "#E63946"

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8),
            Inches(0.55),
            Inches(2.9),
            Inches(0.5),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor.from_string(color[1:])
        badge.line.fill.background()
        tf = badge.text_frame
        tf.text = f"Data Health: {status}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        self._apply_font(tf, size=12, bold=True)

        self._add_slide_notes(slide, f"Data health badge: {status}")

    def _add_slide_notes(self, slide, note: str) -> None:
        notes = slide.notes_slide.notes_text_frame
        if notes.text:
            notes.text += f"\n{note}"
        else:
            notes.text = note

    @staticmethod
    def _parse_number(value: str) -> float | None:
        if value is None:
            return None
        if isinstance(value, (int, float)):
            return float(value)
        text = str(value).replace(",", "").replace("%", "").strip()
        if not text:
            return None
        try:
            return float(text)
        except ValueError:
            return None

    def _format_delta(self, current: str, previous: str) -> str:
        curr = self._parse_number(current)
        prev = self._parse_number(previous)
        if curr is None or prev is None:
            return ""
        diff = curr - prev
        if prev != 0:
            pct = diff / prev
            return f"{diff:,.2f} ({pct:.1%})"
        return f"{diff:,.2f}"

    @staticmethod
    def _wrap_text(text: str, max_len: int) -> str:
        words = text.split()
        lines = []
        current = []
        for word in words:
            if sum(len(w) for w in current) + len(current) + len(word) > max_len:
                lines.append(" ".join(current))
                current = [word]
            else:
                current.append(word)
        if current:
            lines.append(" ".join(current))
        return "\n".join(lines)

    @staticmethod
    def _extract_recommendations(bullets: List[str]) -> List[str]:
        recs = [b for b in bullets if b.lower().startswith("rec:")]
        if recs:
            return [b[4:].strip() for b in recs]
        return []


    def _should_add_mapping_slide(self, mapping: MappingResult) -> bool:
        return len(mapping.mapping) > 10 or len(mapping.new_columns) > 10

    def _add_mapping_detail_slide(self, prs: Presentation, mapping: MappingResult) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header_bar(slide)
        title_box = slide.shapes.add_textbox(self.margin, Inches(0.6), Inches(8.8), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = "Schema Mapping Details"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        self._apply_font(title_tf, size=28, bold=True)

        table = slide.shapes.add_table(
            rows=min(len(mapping.mapping) + 1, 15),
            cols=2,
            left=self.margin,
            top=Inches(1.5),
            width=Inches(9.0),
            height=Inches(5.0),
        ).table
        table.cell(0, 0).text = "Raw Column"
        table.cell(0, 1).text = "Mapped To"
        self._style_table_header(table)

        for idx, (raw, target) in enumerate(mapping.mapping.items(), start=1):
            if idx >= 15:
                break
            raw_label = raw if len(raw) <= 40 else raw[:37].rstrip() + "..."
            target_label = target if len(target) <= 40 else target[:37].rstrip() + "..."
            table.cell(idx, 0).text = raw_label
            table.cell(idx, 1).text = target_label
            self._apply_font(table.cell(idx, 0).text_frame, size=16)
            self._apply_font(table.cell(idx, 1).text_frame, size=16)
